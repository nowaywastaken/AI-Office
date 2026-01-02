from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Optional, List, Tuple
from io import BytesIO


class PPTEngine:
    """
    A comprehensive PowerPoint presentation generator with full formatting control.
    Supports: slides, shapes, text boxes, images, tables, and styling.
    """
    
    def __init__(self):
        self.prs = Presentation()
        self.current_slide = None
        self._set_default_size()
    
    def _set_default_size(self):
        """Set default slide size (16:9)."""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def set_slide_size(self, width: float, height: float, unit: str = 'inches'):
        """Set custom slide size."""
        if unit == 'cm':
            self.prs.slide_width = Cm(width)
            self.prs.slide_height = Cm(height)
        else:
            self.prs.slide_width = Inches(width)
            self.prs.slide_height = Inches(height)
        return self
    
    # ==================== SLIDE OPERATIONS ====================
    
    def add_slide(self, layout_index: int = 6):
        """
        Add a new slide.
        
        Layout indices:
        0 = Title Slide, 1 = Title and Content, 2 = Section Header,
        3 = Two Content, 4 = Comparison, 5 = Title Only, 6 = Blank,
        7 = Content with Caption, 8 = Picture with Caption
        """
        slide_layout = self.prs.slide_layouts[layout_index]
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        return self.current_slide
    
    def add_title_slide(self, title: str, subtitle: str = ''):
        """Add a title slide."""
        slide = self.add_slide(layout_index=0)
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        return slide
    
    def add_content_slide(self, title: str, content: List[str]):
        """Add a title and content slide."""
        slide = self.add_slide(layout_index=1)
        slide.shapes.title.text = title
        
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, item in enumerate(content):
                if i == 0:
                    tf.paragraphs[0].text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
        return slide
    
    # ==================== TEXT BOX OPERATIONS ====================
    
    def add_text_box(
        self,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_name: Optional[str] = None,
        font_size: Optional[float] = None,
        bold: bool = False,
        italic: bool = False,
        color: Optional[str] = None,
        alignment: str = 'left',
        unit: str = 'inches'
    ):
        """Add a text box to the current slide."""
        if self.current_slide is None:
            self.add_slide()
        
        # Convert units
        conv = Cm if unit == 'cm' else Inches
        
        textbox = self.current_slide.shapes.add_textbox(
            conv(left), conv(top), conv(width), conv(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        
        # Alignment
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT
        }
        p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
        
        # Font formatting
        run = p.runs[0]
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = RGBColor.from_string(color.lstrip('#'))
        
        return textbox
    
    # ==================== SHAPE OPERATIONS ====================
    
    def add_shape(
        self,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: Optional[str] = None,
        unit: str = 'inches'
    ):
        """
        Add a shape to the current slide.
        
        Shape types: 'rectangle', 'oval', 'rounded_rectangle', 'triangle'
        """
        if self.current_slide is None:
            self.add_slide()
        
        conv = Cm if unit == 'cm' else Inches
        
        shape_map = {
            'rectangle': MSO_SHAPE.RECTANGLE,
            'oval': MSO_SHAPE.OVAL,
            'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE
        }
        
        shape = self.current_slide.shapes.add_shape(
            shape_map.get(shape_type, MSO_SHAPE.RECTANGLE),
            conv(left), conv(top), conv(width), conv(height)
        )
        
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip('#'))
        
        return shape
    
    # ==================== IMAGE OPERATIONS ====================
    
    def add_image(
        self,
        image_path: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        unit: str = 'inches'
    ):
        """Add an image to the current slide."""
        if self.current_slide is None:
            self.add_slide()
        
        conv = Cm if unit == 'cm' else Inches
        
        kwargs = {'left': conv(left), 'top': conv(top)}
        if width:
            kwargs['width'] = conv(width)
        if height:
            kwargs['height'] = conv(height)
        
        return self.current_slide.shapes.add_picture(image_path, **kwargs)
    
    # ==================== TABLE OPERATIONS ====================
    
    def add_table(
        self,
        data: List[List[str]],
        left: float,
        top: float,
        width: float,
        height: float,
        unit: str = 'inches'
    ):
        """Add a table to the current slide."""
        if self.current_slide is None:
            self.add_slide()
        
        conv = Cm if unit == 'cm' else Inches
        
        rows = len(data)
        cols = len(data[0]) if data else 0
        
        table = self.current_slide.shapes.add_table(
            rows, cols,
            conv(left), conv(top), conv(width), conv(height)
        ).table
        
        for i, row_data in enumerate(data):
            for j, cell_text in enumerate(row_data):
                table.cell(i, j).text = str(cell_text)
        
        return table
    
    # ==================== SAVE OPERATIONS ====================
    
    def save(self, filepath: str):
        """Save presentation to file."""
        self.prs.save(filepath)
        return filepath
    
    def save_to_bytes(self) -> bytes:
        """Save presentation to bytes (for streaming)."""
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def create_blank(self):
        """Return the underlying presentation object."""
        return self.prs
